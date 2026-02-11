"""Generate CatalystNeuro branded PowerPoint slides.

Style inspired by existing CatalystNeuro presentations:
- Clean white backgrounds, minimal decoration
- Bold navy/black titles, simple dash bullets
- Small logo in bottom-right of content slides
- Subtle footer with presenter info
- Greg Dunn neuroscience art backgrounds on title, divider, and closing slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
import urllib.request
import tempfile
from PIL import Image

# --- Brand Colors ---
CN_NAVY = RGBColor(2, 18, 66)          # #021242
CN_BLUE = RGBColor(7, 101, 165)        # #0765A5
CN_LIGHT_BLUE = RGBColor(94, 155, 196) # #5E9BC4
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 246, 248)
MEDIUM_GRAY = RGBColor(150, 150, 160)
BODY_TEXT = RGBColor(50, 50, 55)
BLACK = RGBColor(0, 0, 0)
FOOTER_BG = RGBColor(220, 225, 235)    # #dce1eb light silver-blue
FOOTER_TEXT = RGBColor(31, 73, 125)    # #1f497d dark blue

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets")
LOGO_HORIZONTAL = os.path.join(ASSETS_DIR, "logo_horizontal_light.png")
LOGO_HORIZONTAL_TRANSPARENT = os.path.join(ASSETS_DIR, "logo_horizontal_transparent.png")
LOGO_SQUARE = os.path.join(ASSETS_DIR, "logo_square.png")

# Presentation metadata
PRESENTER = "Benjamin Dichter"
PRES_TITLE = "NWB Data Conversion & DANDI Publishing"

# --- Greg Dunn Art Image Pool ---
# Licensing: gregadunn.com — attribution required on each slide using art
ART_IMAGES = [
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2021/07/36-x-48-cortical-columns-esque-full-painting-1200-lines.jpg",
        "title": "Cortical Columns",
        "year": "2021",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2018/09/Neurogenesis-I-final.jpg",
        "title": "Neurogenesis I",
        "year": "2018",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2012/05/hippocampus-II-small.jpg",
        "title": "Hippocampus II",
        "year": "2012",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2012/05/NG2-flare.jpg",
        "title": "NG2 Flare",
        "year": "2012",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2012/05/Two-Pyramidals-16-X-20.jpg",
        "title": "Two Pyramidals",
        "year": "2012",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2012/05/Beyond-the-Horizon.jpg",
        "title": "Beyond the Horizon",
        "year": "2012",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2016/03/Brainbow-Hippocampus-blue-and-gold.jpg",
        "title": "Brainbow Hippocampus",
        "year": "2016",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2024/02/myelination-ink-painting-2023.jpg",
        "title": "Myelination",
        "year": "2023",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2019/06/Cortical-Circuitboard-purple-2023.jpg",
        "title": "Cortical Circuitboard",
        "year": "2019",
    },
    {
        "url": "https://www.gregadunn.com/wp-content/uploads/2016/10/retina-in-inks-1.jpg",
        "title": "Retina in Inks",
        "year": "2016",
    },
]


def download_image(url):
    """Download an image to a temp file and return the path."""
    suffix = os.path.splitext(url.split("?")[0])[-1] or ".jpg"
    fd, path = tempfile.mkstemp(suffix=suffix)
    os.close(fd)
    urllib.request.urlretrieve(url, path)
    return path


def _set_shape_alpha(shape, alpha_pct):
    """Set fill transparency on a shape using lxml (python-pptx lacks native alpha).

    alpha_pct: 0 = fully transparent, 100 = fully opaque.
    """
    # Access the actual lxml element via _SolidFill._solidFill
    solidFill_elem = shape.fill._fill._solidFill
    srgb = solidFill_elem.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Remove any existing alpha element
    for child in srgb.findall(qn("a:alpha")):
        srgb.remove(child)
    # Add alpha (value is in 1000ths of a percent, so 50% = 50000)
    alpha_elem = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
    srgb.append(alpha_elem)


def add_art_background(slide, art_info, overlay_color=WHITE, overlay_alpha=55):
    """Add a Greg Dunn art image as full-slide background with semi-transparent overlay.

    The image keeps its original aspect ratio and is enlarged to cover the
    entire slide, then cropped to fit (like CSS object-fit: cover).

    overlay_alpha: 0 = fully transparent overlay, 100 = fully opaque.
    """
    img_path = download_image(art_info["url"])

    # Get native aspect ratio
    with Image.open(img_path) as img:
        img_aspect = img.width / img.height
    slide_aspect = SLIDE_WIDTH / SLIDE_HEIGHT

    # Place image at slide size, then crop to preserve aspect ratio
    pic = slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    if img_aspect > slide_aspect:
        # Image is wider — crop left and right equally
        crop = (1 - slide_aspect / img_aspect) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        # Image is taller — crop top and bottom equally
        crop = (1 - img_aspect / slide_aspect) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop

    # Add semi-transparent overlay rectangle
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = overlay_color
    overlay.line.fill.background()
    _set_shape_alpha(overlay, overlay_alpha)

    # Clean up temp file
    try:
        os.unlink(img_path)
    except OSError:
        pass


def add_text_backing(slide, left, top, width, height, color=CN_NAVY, alpha=40):
    """Add a semi-transparent colored rectangle as a text backing for readability."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    _set_shape_alpha(rect, alpha)
    return rect


def add_art_attribution(slide, art_info, above_footer=False):
    """Add attribution text for Greg Dunn art per licensing requirements.

    above_footer: if True, position above the footer bar instead of at the bottom edge.
    """
    text = f"Art: {art_info['title']}, {art_info['year']}, Greg Dunn | gregadunn.com"
    if above_footer:
        top = SLIDE_HEIGHT - Inches(0.44) - Inches(0.35)
    else:
        top = SLIDE_HEIGHT - Inches(0.35)
    txBox = slide.shapes.add_textbox(Inches(0.3), top, Inches(6), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(80, 80, 95)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    return txBox


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             font_color=BODY_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, bullets, font_size=20,
                font_color=BODY_TEXT, spacing=0.65):
    """Add a bulleted list using dash markers, matching the reference style."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(bullets) * spacing + 0.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"-   {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
        p.space_before = Pt(12)
        p.space_after = Pt(4)

    return txBox


def add_footer(slide, section=""):
    """Light footer bar flush to bottom with horizontal logo, presenter, title, section."""
    bar_height = Inches(0.44)
    bar_top = SLIDE_HEIGHT - bar_height

    # Light silver-blue background bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), bar_top, SLIDE_WIDTH, bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BG
    bar.line.fill.background()

    text_y = bar_top + Inches(0.06)

    # Horizontal transparent logo (left, vertically centered in footer)
    logo_width = Inches(1.2)
    logo_pic = slide.shapes.add_picture(
        LOGO_HORIZONTAL_TRANSPARENT, Inches(0.15), bar_top, logo_width
    )
    # Center vertically: offset by (bar_height - logo_height) / 2
    logo_pic.top = bar_top + (bar_height - logo_pic.height) // 2

    # Presenter name
    add_text(
        slide, Inches(2.0), text_y, Inches(3), Inches(0.35),
        PRESENTER, font_size=12, font_color=FOOTER_TEXT, bold=True
    )

    # Presentation title (center)
    add_text(
        slide, Inches(4.5), text_y, Inches(5.3), Inches(0.35),
        PRES_TITLE, font_size=12, font_color=FOOTER_TEXT, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Section (right)
    if section:
        add_text(
            slide, Inches(9.5), text_y, Inches(3.5), Inches(0.35),
            section, font_size=12, font_color=FOOTER_TEXT, bold=True,
            alignment=PP_ALIGN.RIGHT
        )


# -- Slide 1: Title -----------------------------------------------------------

def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title (large, bold, left-aligned)
    add_text(
        slide, Inches(0.8), Inches(0.6), Inches(7.5), Inches(2.5),
        "NWB Data Conversion\n& DANDI Publishing",
        font_size=48, font_color=CN_NAVY, bold=True
    )

    # Subtitle
    add_text(
        slide, Inches(0.8), Inches(3.5), Inches(7), Inches(0.8),
        "Standardizing Neurophysiology Data for Open Science",
        font_size=24, font_color=MEDIUM_GRAY
    )

    # Presenter info
    add_text(
        slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.8),
        "Ben Dichter",
        font_size=28, font_color=BLACK
    )
    add_text(
        slide, Inches(0.8), Inches(5.6), Inches(5), Inches(0.5),
        "Founder, CatalystNeuro",
        font_size=20, font_color=MEDIUM_GRAY
    )

    # Square logo (vertically centered, right side)
    slide.shapes.add_picture(
        LOGO_SQUARE, Inches(9.5), Inches(2.0), Inches(3.0)
    )


# -- Slide 2: Section Divider -------------------------------------------------

def make_section_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)

    # Section label (small)
    add_text(
        slide, Inches(0.8), Inches(2.2), Inches(6), Inches(0.6),
        "The Challenge", font_size=20, font_color=MEDIUM_GRAY
    )

    # Topic title (large)
    add_text(
        slide, Inches(0.8), Inches(3.0), Inches(10), Inches(1.8),
        "Why neurophysiology data\nneeds standardization",
        font_size=42, font_color=CN_NAVY, bold=True
    )

    add_footer(slide, section="The Challenge")


# -- Slide 3: Content (Bullet Points) -----------------------------------------

def make_content_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title (bold, top)
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(11), Inches(1.0),
        "Need for neurophysiology data standard",
        font_size=36, font_color=BLACK, bold=True
    )

    # Bullets
    add_bullets(slide, Inches(0.8), Inches(1.5), Inches(10), [
        "Data are expensive to collect (money, time, animal use)",
        "Sharing neurophysiology data within a lab and with collaborators is tedious",
        "Sharing scientific software is also difficult",
        "It is often easier to just collect new data and build new tools",
        "Difficult for labs to devote time and effort to sharing data and software",
    ], font_size=22, spacing=0.75)

    add_footer(slide, section="The Challenge")


# -- Slide 4: Two-Column Layout -----------------------------------------------

def make_two_column_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(11), Inches(1.0),
        "Our Approach",
        font_size=36, font_color=BLACK, bold=True
    )

    # Left column header
    add_text(
        slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.6),
        "NWB Conversion", font_size=26, font_color=CN_BLUE, bold=True
    )
    add_bullets(slide, Inches(1.0), Inches(2.3), Inches(5), [
        "Automated pipeline from 40+ formats",
        "Full metadata preservation",
        "Validation against NWB schema",
        "Handles multi-modal recordings",
    ], font_size=20, spacing=0.65)

    # Right column header
    add_text(
        slide, Inches(7.0), Inches(1.6), Inches(5), Inches(0.6),
        "DANDI Publishing", font_size=26, font_color=CN_BLUE, bold=True
    )
    add_bullets(slide, Inches(7.2), Inches(2.3), Inches(5), [
        "Upload to DANDI Archive",
        "DOI minting for citation",
        "BIDS-compatible organization",
        "Streaming access via LINDI",
    ], font_size=20, spacing=0.65)

    add_footer(slide, section="Our Approach")


# -- Slide 5: Closing ---------------------------------------------------------

def make_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Horizontal logo centered
    slide.shapes.add_picture(
        LOGO_HORIZONTAL, Inches(3.8), Inches(1.5), Inches(5.5)
    )

    # Thank you
    add_text(
        slide, Inches(1), Inches(3.8), Inches(11.3), Inches(1.0),
        "Thank You",
        font_size=42, font_color=CN_NAVY, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Contact info
    add_text(
        slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
        "ben.dichter@catalystneuro.com",
        font_size=20, font_color=CN_BLUE,
        alignment=PP_ALIGN.CENTER
    )
    add_text(
        slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
        "www.catalystneuro.com  |  github.com/catalystneuro",
        font_size=18, font_color=MEDIUM_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    add_footer(slide)


# -- Main ---------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    make_title_slide(prs)
    make_section_divider(prs)
    make_content_slide(prs)
    make_two_column_slide(prs)
    make_closing_slide(prs)

    output_path = os.path.join(os.path.dirname(__file__), "example_slides.pptx")
    prs.save(output_path)
    print(f"Saved to: {output_path}")


if __name__ == "__main__":
    main()
